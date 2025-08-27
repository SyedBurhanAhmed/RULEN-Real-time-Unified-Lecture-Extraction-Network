from flask import Flask, render_template, request, send_file, redirect, url_for, jsonify, send_file, session
import os
import markdown
from werkzeug.utils import secure_filename
from agents.note_generation_agent import generate_notes_agent
from agents.notes_docx_exporter_agent import save_to_docx
from agents.transcription_agent import transcribe_audio_agent
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.docstore.document import Document
import fitz  # for PDF
from pptx import Presentation  # for PPTX
from docx import Document as DocxDocument  # for DOCX
from langchain_community.vectorstores import FAISS
from langchain_huggingface import HuggingFaceEmbeddings
from langchain.chains import RetrievalQA
from agents.quiz_generator_agent import generate_quiz
from agents.quiz_docx_exporter_agent import save_quiz_to_docx
from agents.exam_assets_generator_agent import generate_assets
from agents.markdown_docx_agent import save_markdown_docx


app = Flask(__name__)
app.secret_key = 'supersecretkey'
UPLOAD_FOLDER = 'uploads'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
DOCX_FOLDER = "output"
# Load vector store components
embedding_model = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
dummy_text = "Initial dummy document"
vectorstore = FAISS.from_texts([dummy_text], embedding_model)

# Immediately remove the dummy document
dummy_id = list(vectorstore.index_to_docstore_id.values())[0]
vectorstore.delete([dummy_id])

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/student', methods=['GET'])
def student_dashboard():
    return render_template('dashboard_student.html')


@app.route("/upload_audio", methods=["POST"])
def upload_audio():
    audio = request.files.get("audio")
    if audio:
        filename = secure_filename(audio.filename)
        path = os.path.join(UPLOAD_FOLDER, filename)
        audio.save(path)
        return jsonify({"filename": filename})
    return jsonify({"error": "No file uploaded"}), 400

@app.route('/submit_audio', methods=['POST'])
def submit_audio():
    filename = request.form.get("filename")
    if not filename:
        return "No filename provided", 400

    filepath = os.path.join(UPLOAD_FOLDER, secure_filename(filename))
    if not os.path.exists(filepath):
        return f"File '{filename}' does not exist on server.", 404

    try:
        transcript, lang = transcribe_audio_agent(filepath)
    except Exception as e:
        return f"Transcription failed: {e}", 500

    try:
        result = generate_notes_agent(transcript, lang)
    except Exception as e:
        return f"Note generation failed: {e}", 500

    try:
        docx_path = save_to_docx(result, filename)
    except Exception as e:
        return f"DOCX export failed: {e}", 500

    # Convert notes from Markdown to HTML for rendering in template
    result["notes_local"] = markdown.markdown(result["notes_local"])
    result["notes_english"] = markdown.markdown(result["notes_english"])
    if result.get("converted_text"):
        result["converted_text"] = markdown.markdown(result["converted_text"])
    session['notes_english'] = result["notes_english"]
    session['docx_path'] = docx_path

    session['language'] = result["language"]
    session['notes_local'] = result["notes_local"]
    session['converted_text'] = result.get("converted_text", "")

    return render_template(
        'results.html',
        result=result,
        docx_link=url_for('download_docx', filename=os.path.basename(docx_path)),
        show_notes=True,
        show_material_upload=False,
        show_quiz_assistant=False,
        status_message="✅ Transcription complete. Notes generated and saved."
    )


@app.route('/teacher', methods=['GET', 'POST'])
def teacher_dashboard():
    return render_template('dashboard_teacher.html')  # Can be extended similarly

@app.route('/download/<filename>')
def download_docx(filename):
    return send_file(os.path.join(DOCX_FOLDER, filename), as_attachment=True)

@app.route('/prepare_quiz_decision', methods=['POST'])
def prepare_quiz_decision():
    decision = request.form['decision']
    if decision == "no":
        return redirect(url_for('/'))  # or a simpler end page
    elif decision == "yes":
        lecture_notes = request.form['lecture_notes']
        splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
        chunks = splitter.split_text(lecture_notes)
        documents = [Document(page_content=chunk) for chunk in chunks]
        vectorstore.add_documents(documents)
        result = {
            "language": session.get("language"),
            "notes_local": session.get("notes_local"),
            "notes_english": session.get("notes_english"),
            "converted_text": session.get("converted_text")
        }

        return render_template(
            'results.html',
            result=result,
            docx_link=url_for('download_docx', filename=os.path.basename(session['docx_path'])),
            show_material_upload=True,
            show_quiz_assistant=False,  # only show assistant if user wants quiz
            show_notes=False  # ⬅️ NEW FLAG to hide notes and other sections
        )
    return None


@app.route('/material_decision', methods=['POST'])
def material_decision():
    decision = request.form['material_decision']
    result = {
        "language": session.get("language"),
        "notes_local": session.get("notes_local"),
        "notes_english": session.get("notes_english"),
        "converted_text": session.get("converted_text")
    }

    return render_template(
        'results.html',
        result=result,
        docx_link=url_for('download_docx', filename=os.path.basename(session['docx_path'])),
        show_material_upload=False,
        show_quiz_assistant=True,  # only show assistant if user wants quiz
        show_notes=False  # ⬅️ NEW FLAG to hide notes and other sections
    )


@app.route('/upload_material', methods=['POST'])
def upload_material():
    file = request.files['materials']
    filename = secure_filename(file.filename)
    filepath = os.path.join(UPLOAD_FOLDER, filename)
    file.save(filepath)

    ext = os.path.splitext(filename)[1].lower()
    content = ""
    try:
        if ext == ".pdf":
            with fitz.open(filepath) as doc:
                content = "\n".join([page.get_text() for page in doc])
        elif ext in [".ppt", ".pptx"]:
            prs = Presentation(filepath)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
        elif ext == ".docx":
            doc = DocxDocument(filepath)
            content = "\n".join([para.text for para in doc.paragraphs])
        elif ext == ".txt":
            with open(filepath, "r", encoding="utf-8") as f:
                content = f.read()
        else:
            return "❌ Unsupported file format.", 400
    except Exception as e:
        return f"❌ Error reading file: {str(e)}", 500

    splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
    chunks = splitter.split_text(content)
    documents = [Document(page_content=chunk) for chunk in chunks]
    # print(f"Total documents in vector store: {len(vectorstore.index_to_docstore)}")
    vectorstore.add_documents(documents)
    result = {
        "language": session.get("language"),
        "notes_local": session.get("notes_local"),
        "notes_english": session.get("notes_english"),
        "converted_text": session.get("converted_text")
    }

    return render_template(
        'results.html',
        result=result,
        docx_link=url_for('download_docx', filename=os.path.basename(session['docx_path'])),
        show_material_upload=False,
        show_quiz_assistant=True
    )

@app.route('/start_quiz_assistant', methods=['POST'])
def start_quiz_assistant():
    question_type = request.form.get("type", "all")
    try:
        quiz_markdown = generate_quiz(vectorstore, question_type)  # returns markdown
    except Exception as e:
        return f"Quiz generation failed: {e}", 500

    # Convert Markdown to HTML for safe rendering
    quiz_html = markdown.markdown(quiz_markdown)

    # Save Markdown to .docx
    filename = save_quiz_to_docx(quiz_markdown, "Quiz Questions and Answers")
    docx_link = url_for('download_docx', filename=filename)

    return render_template(
        "quiz_results.html",
        questions=quiz_html,
        docx_link=docx_link
    )




@app.route('/teacher_submit_audio', methods=['POST'])
def teacher_submit_audio():
    audio = request.files['audio_file']
    filename = secure_filename(audio.filename)
    filepath = os.path.join(UPLOAD_FOLDER, filename)
    os.makedirs(UPLOAD_FOLDER, exist_ok=True)
    audio.save(filepath)

    try:
        # Transcribe audio
        transcript, lang = transcribe_audio_agent(filepath)

        # Save the transcript text to a new file (e.g., in a transcripts folder)
        transcripts_folder = os.path.join(UPLOAD_FOLDER, "transcripts")
        os.makedirs(transcripts_folder, exist_ok=True)
        transcript_filename = os.path.splitext(filename)[0] + "_transcript.txt"
        transcript_path = os.path.join(transcripts_folder, transcript_filename)

        with open(transcript_path, "w", encoding="utf-8") as f:
            f.write(transcript)

        # Store only the transcript file path as a reference in the session
        session['teacher_transcript_file'] = transcript_path
        session['teacher_language'] = lang

    except Exception as e:
        return f"Transcription failed: {e}", 500

    return render_template("teacher_material_prompt.html")  # asks: add slides/material?


@app.route('/teacher_material_decision', methods=['POST'])
def teacher_material_decision():
    decision = request.form['decision']
    if decision == 'yes':
        return render_template("teacher_upload_material.html")
    return render_template("teacher_clo_prompt.html")

@app.route('/teacher_upload_material', methods=['POST'])
def teacher_upload_material():
    file = request.files['materials']
    filename = secure_filename(file.filename)
    filepath = os.path.join(UPLOAD_FOLDER, filename)
    file.save(filepath)

    ext = os.path.splitext(filename)[1].lower()
    content = ""
    try:
        if ext == ".pdf":
            with fitz.open(filepath) as doc:
                content = "\n".join([page.get_text() for page in doc])
        elif ext in [".ppt", ".pptx"]:
            prs = Presentation(filepath)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
        elif ext == ".docx":
            doc = DocxDocument(filepath)
            content = "\n".join([para.text for para in doc.paragraphs])
        elif ext == ".txt":
            with open(filepath, "r", encoding="utf-8") as f:
                content = f.read()
        else:
            return "❌ Unsupported file format.", 400
    except Exception as e:
        return f"❌ Error reading file: {str(e)}", 500

    # Split and embed content
    splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
    chunks = splitter.split_text(content)
    documents = [Document(page_content=chunk) for chunk in chunks]
    vectorstore.add_documents(documents)

    # Mark that material has been uploaded
    session['material_uploaded'] = True

    # Proceed to CLO/PLO decision step
    return render_template("teacher_clo_prompt.html")


@app.route('/teacher_clo_decision', methods=['POST'])
def teacher_clo_decision():
    decision = request.form['decision']

    if decision == 'yes':
        # Show upload page
        return render_template("upload_clos.html")
    else:
        # Proceed with generic CLOs/PLOs
        session['clos_text'] = "Use general CLOs and PLOs for quiz and exam generation."
        return redirect(url_for('generate_exam_assets'))


@app.route('/upload_clos', methods=['POST'])
def upload_clos():
    file = request.files['clos_file']
    filename = secure_filename(file.filename)
    filepath = os.path.join(UPLOAD_FOLDER, filename)
    file.save(filepath)

    ext = os.path.splitext(filename)[1].lower()
    clos_text = ""

    try:
        if ext == ".txt":
            with open(filepath, "r", encoding="utf-8") as f:
                clos_text = f.read()
        elif ext == ".docx":
            doc = DocxDocument(filepath)
            clos_text = "\n".join([para.text for para in doc.paragraphs])
        else:
            return "❌ Unsupported CLO file format. Only .txt and .docx are allowed.", 400
    except Exception as e:
        return f"❌ Error reading CLO file: {str(e)}", 500

    session['clos_text'] = clos_text

    return redirect(url_for('generate_exam_assets'))

@app.route('/generate_exam_assets', methods=['GET', 'POST'])
def generate_exam_assets():
    if request.method == 'GET':
        return render_template("generate_assets_form.html")

    asset_type = request.form['asset_type']  # e.g., "quiz", "mid", "final"
    clos_text = session.get("clos_text", "")
    materials_uploaded = session.get("material_uploaded", False)

    # Instead of storing transcript text in session, session stores transcript file path
    transcript_text = ""
    transcript_file = session.get("teacher_transcript_file", None)
    if transcript_file and os.path.exists(transcript_file):
        with open(transcript_file, "r", encoding="utf-8") as f:
            transcript_text = f.read()

    combined_text = ""
    if transcript_text.strip():
        combined_text += f"**Lecture Transcript:**\n{transcript_text}\n\n"
    print(materials_uploaded)
    if materials_uploaded:
        combined_text += "**Supplementary Materials:**\n(Use the uploaded lecture materials in the vectorstore for factual accuracy.)\n\n"

    if clos_text.strip():
        combined_text += f"**CLOs/PLOs:**\n{clos_text}\n"
    else:
        combined_text += "**CLOs/PLOs:** No CLOs/PLOs provided.\n"

    try:
        generated_text = generate_assets(vectorstore, combined_text, asset_type)
    except Exception as e:
        return f"❌ Exam generation failed: {str(e)}", 500

    filename = f"{asset_type}_generated.docx"
    docx_path = save_markdown_docx(generated_text, filename)
    docx_link = url_for("download_docx", filename=os.path.basename(docx_path))

    return render_template(
        "exam_asset_result.html",
        asset_type=asset_type,
        result_text=generated_text,
        docx_link=docx_link
    )


if __name__ == '__main__':
    # app.run(debug=True)
    app.run()